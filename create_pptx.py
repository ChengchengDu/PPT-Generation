from nltk import text
from numpy.lib.npyio import save
from pptx import Presentation

def add_titile(prs, text="Hello World", slide_num=0):
    title_slide_layout = prs.slide_layouts[slide_num]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = text

def add_abstract(prs, slide_num=1, title_text="abstract", abstract_str=""):
    bullet_slide_layout = prs.slide_layouts[slide_num]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    # 当前slide的标题 采用'主标题-副标题形式' 副标题使用当前章节标题
    # for shape in slide.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title_text
    tf = body_shape.text_frame
    tf.text = abstract_str


def add_slide(prs, slide_num=1, title_text='Adding a Bullet Slide', tf_content=None):
    bullet_slide_layout = prs.slide_layouts[slide_num]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    # 当前slide的标题 采用'主标题-副标题形式' 副标题使用当前章节标题
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title_text
   
    tf = body_shape.text_frame
    tf_num = len(tf_content)
    items = list(tf_content.items())
    for i in range(tf_num):
        tf.text = items[i][0]
        p = tf.add_paragraph()
        p.level = 0

        p.text = items[i][1]
        p = tf.add_paragraph()
        p.level = 1


def add_per_slide(prs, slide_lis, slide_num=1):
    """
    slide: list, including two content
    """
    slide_title = slide_lis[0]["ps_title"]
    bullet_slide_layout = prs.slide_layouts[slide_num]
    slide = prs.slides.add_slide(bullet_slide_layout)  # 新增当前层的幻灯片
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = slide_title
   
    tf = body_shape.text_frame
    tf_num = len(slide_lis)
    for i in range(tf_num):
        kp = "&".join(slide_lis[i]["sent_kp"])
        sentence = slide_lis[i]["cur_sum"]
        # tf.text = kp
        p = tf.add_paragraph()
        p.text = kp
        p.level = 1

        p = tf.add_paragraph()
        p.text = sentence
        p.level = 2

def make_slide(slide_dict, save_path):
    prs = Presentation()     
    paper_titile = slide_dict["paper_title"]
    # print("paper_titile")
    # print(paper_titile)
    save_path = save_path + paper_titile + '.pptx'
    add_titile(prs, text=paper_titile)
    add_abstract(prs, abstract_str=slide_dict["paper_abstract"])
    all_slides = slide_dict["slides"]
    for slide in all_slides:
        add_per_slide(prs, slide)
    
    prs.save(save_path)


# if __name__ == "__main__":
#     prs = Presentation()    
#     slide = [
#         {'ps_title': 'Introduction', 'sent_kp': ['vcd differs', 'standard kl'], 'cur_sum': 'The VCD differs from the standard KL in that it captures a notion of discrepancy between the improved and the initial variational distributions.'},
#         {'ps_title': 'Introduction', 'sent_kp': ['vcd objective', 'predictive performance'], 'cur_sum': 'We show experimentally that the models fitted with the VCD objective have better predictive performance than when we use alternative approaches, including standard VI and the method of Hoffman (2017).'}
#         ]
#     add_titile(prs, text="hello")

#     add_per_slide(prs, slide)
    
#     prs.save('./test.pptx')