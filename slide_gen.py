import json
import re
import nltk
import os
import sys
import numpy as np
import torch
from functools import lru_cache
from nltk.sem.evaluate import Model
from nltk.tokenize import sent_tokenize
from nltk import data, inference
from numpy.lib.npyio import save
from numpy.lib.utils import _set_function_name
from pickle import NEWFALSE

from yaml import tokens
from sklearn.cluster import KMeans
from sklearn.metrics import pairwise_distances_argmin_min
from sentence_transformers import SentenceTransformer
sys.path.append("/home/jazhan/code/document2slides-main/d2s_model/")
sys.path.append("/home/jazhan/code/PacSum-master/code")
from extractor import PacSumExtractorWithBert, PacSumExtractorWithTfIdf
from inference_data import InferDataset
from infer import extract_summary_pac
from extract_keys import WordEmbeddings, SentEmbeddings, SIFRank_plus
from create_pptx import make_slide
from abstract_generate import generate_summary
from transformers import PegasusForConditionalGeneration, PegasusTokenizer
os.environ["TOKENIZERS_PARALLELISM"] = "false"
from nltk import text
from numpy.lib.npyio import save
from pptx import Presentation

class MakeSlide():
    def __init__(self) -> None:
        self.prs = Presentation()
       
    def add_titile(self, text="Hello World", slide_num=0):
        title_slide_layout = self.prs.slide_layouts[slide_num]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = text

    def add_abstract(self, slide_num=1, title_text="abstract", abstract_str=""):
        bullet_slide_layout = self.prs.slide_layouts[slide_num]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = title_text
        tf = body_shape.text_frame
        tf.text = abstract_str

    def add_slide(self, slide_num=1, title_text='Adding a Bullet Slide', tf_content=None):
        bullet_slide_layout = self.prs.slide_layouts[slide_num]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        # 当前slide的标题 采用'主标题-副标题形式' 副标题使用当前章节标题
        for shape in slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = title_text
    
        tf = body_shape.text_frame
        tf_num = len(tf_content)
        items = list(tf_content.items())
        for i in range(tf_num):
            tf.text = items[i][0]
            p = tf.add_paragraph()
            p.level = 0

            p.text = items[i][1]
            p = tf.add_paragraph()
            p.level = 1

    def add_per_slide(self, slide_layer, slide_num=1):
        """
        slide: list, including two content
         # {"ps_title": ps_titl, "slide_cont": slide_cont, "flag": 1}
         slide_cont:
        """
        slide_title = slide_layer["ps_title"]
        bullet_slide_layout = self.prs.slide_layouts[slide_num]
        slide = self.prs.slides.add_slide(bullet_slide_layout)  # 新增当前层的幻灯片
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = slide_title
    
        tf = body_shape.text_frame
        slide_cont = slide_layer["slide_cont"]
        tf_num = len(slide_cont)
        for i in range(tf_num):
            if slide_layer["flag"] == 1:
                kp = slide_cont[i]["kp"]
                cur_sum = slide_cont[i]["cur_sum"]
                {"sent_kp": kp, "cur_sum": cur_sum}
                p = tf.add_paragraph()
                p.text = kp
                p.level = 1

                p = tf.add_paragraph()
                p.text = cur_sum
                p.level = 2
            else:
                cur_sum = slide_cont[i]  # 直接就是一个点
                {"sent_kp": kp, "cur_sum": cur_sum}
                p = tf.add_paragraph()
                p.text = cur_sum
                p.level = 1

                p = tf.add_paragraph()
                p.text = cur_sum
                p.level = 2
           
    def save_slide(self, save_path):
        self.prs.save(save_path)

    def make_slide(self, slide_dict, save_path):    
        paper_titile = slide_dict["paper_title"]
        save_path = save_path + paper_titile + '.pptx'
        self.add_titile(text=paper_titile)
        self.add_abstract(abstract_str=slide_dict["paper_abstract"])
        all_slides = slide_dict["slides"]
        for slide in all_slides:
            self.add_per_slide(slide)      
       

        

class AbstractModel:
    """
    压缩句子
    """
    def __init__(self, absmodel_name='google/pegasus-cnn_dailymail'):
        self.device = 'cuda' if torch.cuda.is_available() else 'cpu'
        self.tokenizer = PegasusTokenizer.from_pretrained(absmodel_name)
        self.model = PegasusForConditionalGeneration.from_pretrained(absmodel_name).to(self.device)
        
    def inference(self, src_text, max_length=200, min_length=20, num_beams=6, length_panalty=2.0, no_repeate_ngram_size=3):

        batch = self.tokenizer(src_text, truncation=True, padding='longest', return_tensors="pt").to(self.device)
        translated = self.model.generate(
            **batch, 
            max_length=max_length, 
            min_length=min_length,
            num_beams=num_beams, 
            length_penalty=length_panalty, 
            no_repeat_ngram_size=no_repeate_ngram_size,
        )
        tgt_text = self.tokenizer.batch_decode(translated, skip_special_tokens=True)
        tgt_text = tgt_text[0].replace("We", "").strip().split('<n>')
        tgt_text = " ".join(tgt_text)

        return tgt_text
            

def filter_sent(sent):
    sent = re.sub(r'[(].*[)]','',sent)
    return sent

def cluster_summary(content_string, model, n_clusters=4):
    sentences = sent_tokenize(content_string, language='english')
    # You would need to download pre-trained models first
    content_embedding = model.encode(sentences)
    # n_clusters = int(np.ceil(len(content_embedding)**0.5))
    if n_clusters > len(sentences):
        n_clusters = len(sentences)
    
    kmeans = KMeans(n_clusters=n_clusters)
    kmeans = kmeans.fit(content_embedding)

    avg = []
    for j in range(n_clusters):
        idx = np.where(kmeans.labels_ == j)[0]
        avg.append(np.mean(idx))
    closest, _ = pairwise_distances_argmin_min(kmeans.cluster_centers_, content_embedding)
    ordering = sorted(range(n_clusters), key=lambda k: avg[k])
    summary = [sentences[closest[idx]] for idx in ordering]

    return summary


class SlideGeneration:
    def __init__(self, model_name='distilbert-base-nli-mean-tokens'):
        sentence_model = SentenceTransformer(model_name)
        self.BERT = WordEmbeddings(model=sentence_model)
        self.SIF = SentEmbeddings(self.BERT, lamda=1.0)
        beta = 0
        lambda1 = 0
        lambda2 = 1.0
        bert_vocab_file = "/home/jazhan/code/PacSum-master/pacssum_models/vocab.txt"
        bert_model_file = "/home/jazhan/code/PacSum-master/pacssum_models/pytorch_model_finetuned.bin"
        bert_config_file = "/home/jazhan/code/PacSum-master/pacssum_models/bert_config.json"
        self.ExtractModel = PacSumExtractorWithBert(
                bert_model_file=bert_model_file,
                bert_config_file=bert_config_file,
                beta=beta,
                lambda1=lambda1,
                lambda2=lambda2,
                )
        self.inference_dataset = InferDataset(vocab_file=bert_vocab_file)

        self.sentence_model = SentenceTransformer()
        self.abs_model = AbstractModel()
        self.slide_dict = {
            "paper_title": "",
            "paper_abstract": "",
            "slides":[],
        }
    
    def extract_summary(self, article):
        article_lis = nltk.sent_tokenize(article)
        extract_num = len(article_lis) // 4
        inference_datasets_iterator = self.inference_dataset._load_once_doc_bert(article_lis)
        summaries = self.ExtractModel.inference_extract_summary(inference_datasets_iterator, extract_num=extract_num)
        return summaries

    def abstract_sum(self, src_text, max_length=200, min_length=20, num_beams=6, length_panalty=2.0, no_repeate_ngram_size=3):
        summaries = self.abs_model.inference(
            src_text, 
            max_length=max_length, 
            min_length=min_length, 
            num_beams=num_beams, 
            length_panalty=length_panalty, 
            no_repeate_ngram_size=no_repeate_ngram_size
        )
        return summaries
    
    def parall_content(self, content_string, ps_titl):
        all_tokens = len(content_string.split())
        max_length, min_length = all_tokens // 2, all_tokens // 4
        content_string = self.abstract_sum(content_string, max_length=max_length, min_length=min_length)
        content_lis = nltk.sent_tokenize(content_string)
        part_slide = {"ps_title": ps_titl, "sent_kp": "no", "content": content_lis, "flag":0}
        self.slide_dict["slides"].append(part_slide)   # 插入一页幻灯片， flag=0表示没有使用关键词， 需要将content中每一个句子作为一个要点


    def slide_gen(self, json_path='/home/jazhan/code/document2slides-main/input/ACL_paper_json/965.json'):
        with open(json_path, 'r', encoding="utf-8") as f:
            data = json.load(f)   
        extractive_summary_dict = {
            "title": "",
            "abstract": [],
            "extractive_summary":[],
        }
        self.slide_dict["paper_title"] = data["title"]
        abstract = self.abstract_sum(data["abstract"])
        self.slide_dict["paper_abstract"] = abstract
        num2sec = {}
        section_names = []
        for section_text in data["text"]:
            section_name = section_text["section"]
            section_names.append(section_name)
            print("section name", section_name)
            content_string = section_text['strings']
            print("content_string", content_string)
            if content_string.startswith(section_name):
                content_string = content_string.replace(section_name, '').strip()
            if section_name.lower() == "related work":
                continue
            n = section_text['n']   # 章节编号
            if len(n.split('.')) <= 2:
                n = str(float(n))
            num_sec = n.split('.')[0]
            if num_sec not in num2sec:
                num2sec[num_sec] = section_name
            if len(nltk.sent_tokenize(content_string)) < 4:
                # 直接作为若干个要点作为一个PPT 这些大部分是总领性的文字
                content_lis = nltk.sent_tokenize(content_string)
                part_slide = {"ps_title": section_name, "slide_cont": content_lis, "flag":0}
                self.slide_dict["slides"].append(part_slide)
                continue
            if n.split('.')[1] == '0':
                # 说明不是小章节
                # introduction的部分如何设计？直接生成式摘要，基本是每四个句子生成一个句子
                tokens_all = len(content_string.split())
                max_length, min_length = tokens_all // 4, tokens_all // 8
                content_string = self.abstract_sum(content_string, max_length=max_length, min_length=min_length)
                content_lis = nltk.sent_tokenize(content_string)
                for i in range(0, len(content_lis), 4):
                    end = i + 4 if (i + 4) < len(content_lis) else len(content_lis)
                    cur_sum = content_lis[i: end]
                    part_slide = {"ps_title": section_name, "slide_cont": cur_sum, "flag":0}   # 每一页幻灯片，slide_cont是每张幻灯片的内容, flag=0表示没有关键词
                    self.slide_dict["slides"].append(part_slide)
                continue

            # 获取当前小章节幻灯片的大标题
            if len(n.split('.')) > 2 or (len(n.split('.')) == 2 and n.split('.')[1] != '0'):
                ps_titl = num2sec[num_sec] + "--" + section_name
            else:     
                ps_titl = num2sec[num_sec]   
            ext_summary = self.extract_summary(content_string)
            for sum in ext_summary:
                if len(sum.split()) < 10:
                    ext_summary.remove(sum)
            slide_cont = []              
            for i in range(len(ext_summary)):
                cur_sum = ext_summary[i]
                sent_kp_tu, sent_trkp, sent_nn_candi_embed = SIFRank_plus(cur_sum, self.SIF, N=2, openie=False)
                if len(cur_sum.split()) > 50:
                    print("abstract cur sum {}".format(len(cur_sum.split())))
                    cur_sum = self.abstract_sum(cur_sum, max_length=25, min_length=10)
                sent_kp = [tp[0] for tp in sent_kp_tu if (len(tp[0].split()) > 1 and tp[0] not in ps_titl)]
                sent_kp = ' and '.join(sent_kp)
                slide_cont.append({"sent_kp": sent_kp, "cur_sum": cur_sum})
                if len(slide_cont) >= 2:
                    part_slide = {"ps_title": ps_titl, "slide_cont": slide_cont, "flag": 1}
                    self.slide_dict["slides"].append(part_slide)
                    slide_cont = []
                
        save_path = "/home/jazhan/code/document2slides-main/d2s_model/demo_pptx/"
        print("slide dict")
        print(self.slide_dict)
        make_slide(self.slide_dict, save_path=save_path)


if __name__ == "__main__":
    slide_gen_obj = SlideGeneration()
    slide_gen_obj.slide_gen(json_path='/home/jazhan/code/document2slides-main/input/ACL_paper_json/965.json')
    print("slide generation has done")


